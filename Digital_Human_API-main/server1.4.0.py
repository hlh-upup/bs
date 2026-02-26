"""
    添加可选择对指定页数插入音频
    修复了：第二次上传自定义音频没有清空上一次的音频
            数字人拼接没有贴边
            数字人太小
    增加了初始化文件夹（删除文件）
"""
import sys, os
base_dir = os.path.dirname(__file__)
if base_dir not in sys.path:
    sys.path.insert(0, base_dir)
# 移除其它同名旧工程路径避免干扰
_remove = []
for p in list(sys.path):
    if p.startswith('F:\\human\\Digital_Human_API') and p != base_dir:
        _remove.append(p)
for p in _remove:
    try:
        sys.path.remove(p)
    except ValueError:
        pass

import base64
import json
import os
import logging
import time
from flask import Flask, request, jsonify, g, send_file
from flask_cors import CORS
import functools
import traceback
import wave
import struct
import shutil
import io

import numpy as np
import cv2
try:
    # AnimeGANv2 推理封装（可选）
    from util.anime_cartoonizer import (
        cartoonize_animegan_onnx,
        resolve_weight_path,
        is_session_cached,
        get_session_providers,
    )
except Exception as _anime_e:
    cartoonize_animegan_onnx = None  # type: ignore
    resolve_weight_path = None  # type: ignore
    is_session_cached = None  # type: ignore
    get_session_providers = None  # type: ignore

# WBC (White-box Cartoonization) via TensorFlow SavedModel (HuggingFace)
try:
    from util.wbc_cartoonizer import wbc_cartoonize, is_wbc_model_loaded  # type: ignore
except Exception as _wbc_e:
    wbc_cartoonize = None  # type: ignore
    is_wbc_model_loaded = lambda: False  # type: ignore

from Main import *
from util.Function import Verification, Encode_Video, Write_Json

import concurrent.futures

# 创建线程池
executor = concurrent.futures.ThreadPoolExecutor(max_workers=5)

def create_blank_video(output_path, duration=10.0):
    """创建空白视频文件"""
    try:
        from moviepy.editor import ImageClip, ColorClip
        from moviepy.video.io.ffmpeg_tools import ffmpeg_extract_subclip

        # 创建空白视频 (1920x1080, 白色背景)
        blank_clip = ColorClip((1920, 1080), color=(255, 255, 255), duration=duration)
        blank_clip.write_videofile(output_path, fps=24, codec='libx264', audio=False)
        blank_clip.close()
        return True
    except Exception as e:
        logger.error(f"创建空白视频失败: {e}")
        return False

############################################
# 全局日志配置
############################################
LOG_FORMAT = "%(asctime)s | %(levelname)s | %(name)s | %(message)s"
LOG_LEVEL = os.environ.get("APP_LOG_LEVEL", "INFO").upper()
RAISE_ON_400 = os.environ.get("APP_RAISE_ON_400", "0") in ("1","true","True")
LOG_ERROR_BODY = os.environ.get("APP_LOG_ERROR_BODY", "1") in ("1","true","True")
DEFAULT_MODE = os.environ.get("APP_DEFAULT_MODE", "").strip()  # 缺省Mode自动回退，可设为 Test / VITS / ...
logging.basicConfig(level=LOG_LEVEL, format=LOG_FORMAT)
logger = logging.getLogger("server")

app = Flask(__name__)
CORS(app)  # 允许跨域请求

############################################
# 通用异常装饰器（推理端点使用）
############################################
def log_exceptions(re_raise: bool = False):
    """装饰路由函数，捕获异常写入日志；可通过 APP_RAISE_ON_INFERENCE 环境变量或参数 re_raise 决定是否向上抛出。
    若不抛出，则返回 500 JSON，并带上 error 字段。
    可选 APP_INFERENCE_TRACE=1 时在 JSON 中附加 stack 字段（截断）。"""
    def decorator(fn):
        @functools.wraps(fn)
        def wrapper(*args, **kwargs):
            try:
                return fn(*args, **kwargs)
            except Exception as e:
                logger.exception(f"Inference endpoint {fn.__name__} failed: {e}")
                should_raise = re_raise or os.environ.get("APP_RAISE_ON_INFERENCE", "0") in ("1","true","True")
                if should_raise:
                    raise
                resp = {"result": "Failed", "error": str(e)}
                if os.environ.get("APP_INFERENCE_TRACE", "0") in ("1","true","True"):
                    tb = traceback.format_exc()
                    if len(tb) > 1200:
                        tb = tb[:1200] + "...<truncated>"
                    resp["stack"] = tb
                return jsonify(resp), 500
        return wrapper
    return decorator

# 记录所有请求 & 响应耗时
@app.before_request
def _log_request_start():
    g._start_time = time.time()
    logger.info(f"REQUEST {request.method} {request.path} ip={request.remote_addr} args={dict(request.args)}")
    if LOG_LEVEL == 'DEBUG':
        # 头信息
        logger.debug(f"HEADERS {dict(request.headers)}")
        # 表单字段
        if request.form:
            logger.debug(f"FORM {request.form.to_dict()}")
        # 文件信息（仅列出名字和大小）
        if request.files:
            file_meta = {k: {"filename": v.filename, "content_length": request.content_length} for k, v in request.files.items()}
            logger.debug(f"FILES {file_meta}")
    if request.method in ("POST","PUT","PATCH"):
        try:
            if request.is_json:
                logger.info(f"JSON_BODY {request.get_json(silent=True)}")
        except Exception as e:
            logger.warning(f"读取请求JSON失败: {e}")

@app.after_request
def _log_request_end(resp):
    try:
        cost = (time.time() - getattr(g, '_start_time', time.time())) * 1000
        logger.info(f"RESPONSE {request.method} {request.path} status={resp.status_code} cost_ms={cost:.2f}")
        if LOG_ERROR_BODY and resp.status_code >= 400:
            try:
                body_preview = resp.get_data(as_text=True)
                if len(body_preview) > 800:
                    body_preview = body_preview[:800] + '...<truncated>'
                logger.error(f"ERROR_BODY {request.method} {request.path} -> {body_preview}")
            except Exception as log_e:
                logger.warning(f"读取响应体失败: {log_e}")
    except Exception as e:
        logger.warning(f"after_request 日志异常: {e}")
    return resp

# 全局异常兜底
@app.errorhandler(Exception)
def _handle_exception(e):
    logger.exception(f"UNCAUGHT_EXCEPTION path={request.path} error={e}")
    return jsonify(result="Failed", error=str(e)), 500

# 404 调试增强
@app.errorhandler(404)
def _handle_404(e):
    routes = sorted([r.rule for r in app.url_map.iter_rules()])
    logger.warning(f"404 Not Found: {request.path}. 可用路由: {routes}")
    return jsonify(result="NotFound", path=request.path, routes=routes), 404

############################################
# 图片工具 & 二次元化实现（轻量版：OpenCV）
############################################

def _b64_to_image(img_b64: str):
    try:
        data = base64.b64decode(img_b64)
        arr = np.frombuffer(data, np.uint8)
        img = cv2.imdecode(arr, cv2.IMREAD_COLOR)
        return img
    except Exception as e:
        raise RuntimeError(f"base64解码失败: {e}")

def _image_to_b64(img) -> str:
    ok, buf = cv2.imencode('.png', img)
    if not ok:
        raise RuntimeError("图像编码PNG失败")
    return base64.b64encode(buf.tobytes()).decode('utf-8')

def cartoonize_stylization(img, sigma_s=60, sigma_r=0.07):
    # OpenCV 自带卡通化风格（非真实Anime，但更二次元化）
    try:
        cartoon = cv2.stylization(img, sigma_s=int(sigma_s), sigma_r=float(sigma_r))
        return cartoon
    except Exception as e:
        raise RuntimeError(f"stylization失败: {e}")

def cartoonize_bilateral(img, bilateral_times=5, d=9, sigmaColor=150, sigmaSpace=150, edge_block_size=9, edge_C=2):
    # 经典卡通化：多次双边 + 边缘二值化
    color = img.copy()
    for _ in range(int(max(1, bilateral_times))):
        color = cv2.bilateralFilter(color, d=int(d), sigmaColor=int(sigmaColor), sigmaSpace=int(sigmaSpace))

    gray = cv2.cvtColor(img, cv2.COLOR_BGR2GRAY)
    gray = cv2.medianBlur(gray, 7)
    edges = cv2.adaptiveThreshold(gray, 255, cv2.ADAPTIVE_THRESH_MEAN_C, cv2.THRESH_BINARY, int(edge_block_size) | 1, int(edge_C))
    edges = cv2.cvtColor(edges, cv2.COLOR_GRAY2BGR)
    cartoon = cv2.bitwise_and(color, edges)
    return cartoon

@app.route('/Cartoonize_Image', methods=['POST'])
def Cartoonize_Image():
    """二次元化（多模式，无自动降级）。
    请求: { 
        User, 
        Img(base64-raw), 
        Mode: 'animegan_v2'|'wbc'|'cv_stylize'|'bilateral', 
        Style?: 'hayao'|'shinkai'|'paprika'|'celeba' (仅 animegan_v2), 
        Params?: { max_side?: number, sigma_s?: number, sigma_r?: number, bilateral_times?:int, d?:int, sigmaColor?:int, sigmaSpace?:int, edge_block_size?:int, edge_C?:int }
    }
    响应: { result: 'Success', Img, mode_used, style_used?, debug:{...} }
    说明：不进行兜底，前端选择哪种模式就按哪种执行。
    """
    try:
        POST_JSON = request.get_json() or {}
        user = POST_JSON.get('User')
        img_b64 = POST_JSON.get('Img')
        requested_mode = (POST_JSON.get('Mode') or 'animegan_v2').lower()
        style = (POST_JSON.get('Style') or 'hayao').lower()
        params = POST_JSON.get('Params') or {}
        

        if not img_b64:
            return jsonify(result="Failed", error="缺少图片数据")

        img = _b64_to_image(img_b64)
        if img is None:
            return jsonify(result="Failed", error="图片解码失败")

        h, w = img.shape[:2]
        if h == 0 or w == 0:
            return jsonify(result="Failed", error="空图像")

        max_side = int(params.get('max_side', 1600))
        scale = 1.0
        if max(h, w) > max_side:
            scale = max_side / float(max(h, w))
            img_proc = cv2.resize(img, (int(w * scale), int(h * scale)), interpolation=cv2.INTER_AREA)
        else:
            img_proc = img

        providers = None
        weight_path = None
        cached_before = None

        logger.info(f"[Cartoonize] start mode={requested_mode} style={style} size={w}x{h} scale={scale:.4f}")

        t0 = time.time()
        if requested_mode == 'animegan_v2':
            if cartoonize_animegan_onnx is None or resolve_weight_path is None:
                return jsonify(result="Failed", error="AnimeGANv2 模块未加载或环境缺失(onnxruntime)")
            # 记录权重与缓存状态
            weight_path = resolve_weight_path(style)
            try:
                if is_session_cached is not None:
                    cached_before = bool(is_session_cached(style))
            except Exception:
                cached_before = False
            out = cartoonize_animegan_onnx(img_proc, style=style)
            try:
                if get_session_providers is not None:
                    providers = get_session_providers(style)
            except Exception:
                providers = None
        elif requested_mode == 'wbc':
            if wbc_cartoonize is None:
                return jsonify(result="Failed", error="WBC 模块不可用，请确认已安装 tensorflow 与 huggingface_hub" )
            # 注意：WBC 内部自带 resize/crop 到 8 的倍数，先不缩放，由内部处理，再出图后按需还原原尺寸
            infer = wbc_cartoonize(img)
            out = infer["output"]
            providers = "tensorflow"
        elif requested_mode == 'cv_stylize':
            sigma_s = int(params.get('sigma_s', 60))
            sigma_r = float(params.get('sigma_r', 0.07))
            out = cartoonize_stylization(img_proc, sigma_s=sigma_s, sigma_r=sigma_r)
        elif requested_mode == 'bilateral':
            out = cartoonize_bilateral(
                img_proc,
                bilateral_times=int(params.get('bilateral_times', 5)),
                d=int(params.get('d', 9)),
                sigmaColor=int(params.get('sigmaColor', 150)),
                sigmaSpace=int(params.get('sigmaSpace', 150)),
                edge_block_size=int(params.get('edge_block_size', 9)),
                edge_C=int(params.get('edge_C', 2)),
            )
        else:
            return jsonify(result="Failed", error=f"不支持的 Mode: {requested_mode}")

        cost_ms = (time.time() - t0) * 1000.0

        # 恢复原尺寸（若前面缩小过）
        if scale != 1.0 and requested_mode != 'wbc':
            out = cv2.resize(out, (w, h), interpolation=cv2.INTER_CUBIC)
        # 对于 wbc，为了不二次插值，这里使用等比拉伸到原尺寸
        if requested_mode == 'wbc' and (out.shape[0] != h or out.shape[1] != w):
            out = cv2.resize(out, (w, h), interpolation=cv2.INTER_CUBIC)

        out_b64 = _image_to_b64(out)

        # 保存到用户目录（可选）
        if user:
            try:
                _, _, _, save_user_path, _, _ = Create_File(str(user))
                out_path = os.path.join(save_user_path, 'Image_Cartoon.png')
                cv2.imwrite(out_path, out)
            except Exception as se:
                logger.warning(f"保存卡通化图片失败(忽略): {se}")

        logger.info(
            f"[Cartoonize] success mode={requested_mode} style={style if requested_mode=='animegan_v2' else '-'} providers={providers} cost_ms={cost_ms:.2f} out_size={out.shape[1]}x{out.shape[0]}"
        )

        return jsonify(
            result="Success",
            Img=out_b64,
            mode_used=requested_mode,
            style_used=style if requested_mode=='animegan_v2' else None,
            debug={
                "weight_path": weight_path,
                "cached_before": cached_before,
                "providers": providers,
                "cost_ms": cost_ms,
                "in_size": [w, h],
                "out_size": [out.shape[1], out.shape[0]],
                "scale": scale,
                "params": params,
                **({"wbc": infer["debug"]} if requested_mode == 'wbc' else {}),
            },
        )
    except Exception as e:
        logger.exception(f"Cartoonize_Image 失败: {e}")
        return jsonify(result="Failed", error=str(e))

# 列出所有路由
@app.route('/Routes', methods=['GET'])
def list_routes():
    routes = []
    for r in app.url_map.iter_rules():
        if r.rule.startswith('/static'):  # 忽略静态
            continue
        routes.append({"rule": r.rule, "methods": sorted(m for m in r.methods if m not in ['HEAD','OPTIONS'])})
    return jsonify(result=routes)

@app.route('/Get_Inference', methods=['POST'])
def Get_Inference():
    """主推理接口：VITS + SadTalker/Wav2Lip 异步推理
    兼容C#和Vue前端，返回Audio_Video_Inference任务状态
    """
    POST_JSON = request.get_json() or {}
    user = POST_JSON.get('User')

    logger.info(f"Get_Inference: 收到请求，用户: {user}")
    print(f"[GET_INFERENCE_START] 收到推理请求，用户: {user}")

    if not user:
        logger.warning("Get_Inference: 缺少用户参数")
        response_data = {"result": "Failed", "error": "缺少用户参数"}
        logger.info(f"Get_Inference: 参数错误返回数据: {response_data}")
        print(f"[GET_INFERENCE_PARAM_ERROR] 参数错误返回数据: {response_data}")
        return jsonify(response_data)

    try:
        logger.info(f"Get_Inference: 开始为用户 {user} 创建目录结构...")
        # 创建用户目录结构
        result_vits_user_path, result_sadtalker_user_path, result_wav2lip_user_path, save_user_path, _, _ = Create_File(user)
        logger.info(f"Get_Inference: 目录创建成功，保存路径: {save_user_path}")

        # 检查音频文件存在情况，决定推理方式
        ref_wav_path = os.path.join(save_user_path, "Ref_Wav.wav")
        user_wav_path = os.path.join(save_user_path, "User_Wav.wav")

        logger.info(f"Get_Inference: 检查音频文件 - Ref_Wav: {os.path.exists(ref_wav_path)}, User_Wav: {os.path.exists(user_wav_path)}")

        # 设置任务状态为进行中
        logger.info(f"Get_Inference: 开始设置任务状态为进行中...")
        Task_State(save_user_path, "Audio_Video_Inference", "Processing")
        logger.info(f"Get_Inference: 任务状态设置完成")

        # 根据音频文件选择推理方式
        if os.path.exists(user_wav_path):
            # 用户训练的音频 + Wav2Lip推理
            logger.info(f"Get_Inference: 启动用户音频 + Wav2Lip推理，用户: {user}")
            executor.submit(User_Wav_Wav2Lip_Inference, result_wav2lip_user_path, save_user_path)
        elif os.path.exists(ref_wav_path):
            # 参考音频 + VITS + SadTalker推理
            logger.info(f"Get_Inference: 启动VITS + SadTalker推理，用户: {user}")
            executor.submit(VITS_Sadtalker_Inference, result_vits_user_path, result_sadtalker_user_path, save_user_path)
        else:
            # 默认VITS + SadTalker推理
            logger.info(f"Get_Inference: 启动默认VITS + SadTalker推理，用户: {user}")
            executor.submit(VITS_Sadtalker_Inference, result_vits_user_path, result_sadtalker_user_path, save_user_path)

        logger.info(f"Get_Inference: 任务已提交，返回 Audio_Video_Inference")
        # 立即返回任务ID，前端开始轮询
        response_data = {"result": "Audio_Video_Inference"}
        logger.info(f"Get_Inference: 返回数据: {response_data}")
        print(f"[GET_INFERENCE_RETURN] 返回数据: {response_data}")
        return jsonify(response_data)

    except Exception as e:
        logger.exception(f"Get_Inference: 处理失败，用户: {user}, 错误: {e}")
        import traceback
        traceback.print_exc()
        response_data = {"result": "Failed", "error": str(e)}
        logger.info(f"Get_Inference: 异常返回数据: {response_data}")
        print(f"[GET_INFERENCE_ERROR] 异常返回数据: {response_data}")
        return jsonify(response_data)

@app.route('/Login', methods=['POST'])
def Login():
    POST_JSON = request.get_json()
    user_name = POST_JSON.get("User")
    user_password = POST_JSON.get("Password")
    
    try:
        if(Verification(str(user_name), str(user_password))):
            user_result_vits_path, user_result_sadtalker_path, user_result_wav2lip_path, user_data_save_path, _, _ = Create_File(str(user_name))
            Init_File(user_result_vits_path, user_result_sadtalker_path, user_result_wav2lip_path, user_data_save_path)
            return jsonify(result="Success")
        else:
            return jsonify(result="Failed")
    except Exception as e:
        print("Login error:", e)
        return jsonify(result="Failed")
    
#获取状态
@app.route('/Get_State', methods=['POST'])
def Get_State():
    POST_JSON = request.get_json()
    user = POST_JSON.get('User')
    task = POST_JSON.get('Task')

    try:
        # 直接构造状态文件路径，避免创建目录
        user_data_save_path = os.path.normpath(os.path.join('Data', str(user)))
        state_file_path = os.path.join(user_data_save_path, "State.json")

        # 检查状态文件是否存在
        if not os.path.exists(state_file_path):
            logger.warning(f"Get_State: 状态文件不存在 {state_file_path}")
            return jsonify(result="Failed")

        # 直接读取状态文件
        with open(state_file_path, 'r', encoding='utf-8') as f:
            state_data = json.load(f)

        task_state = state_data.get(str(task), "Failed")
        logger.info(f"Get_State: 用户 {user} 任务 {task} 状态 {task_state}")

        return jsonify(result=task_state)
    except Exception as e:
        logger.error(f"Get_State: 查询状态失败，用户 {user} 任务 {task} 错误 {e}")
        return jsonify(result="Failed")
    
# 上传PPT文件并解析备注
@app.route('/Upload_PPT_Parse_Remakes', methods=['POST'])
def Upload_PPT_Parse_Remakes():
    """上传PPT文件并自动解析备注内容，返回解析后的备注JSON"""
    try:
        # 获取上传的PPT文件
        if 'File' not in request.files:
            return jsonify(result="Failed", error="没有找到PPT文件")

        ppt_file = request.files['File']
        if ppt_file.filename == '':
            return jsonify(result="Failed", error="没有选择PPT文件")

        # 获取用户信息
        json_data = request.form.get('Json')
        if not json_data:
            return jsonify(result="Failed", error="缺少用户信息")

        user_info = json.loads(json_data)
        user = user_info.get('User')
        if not user:
            return jsonify(result="Failed", error="缺少用户名")

        _, _, _, save_user_path, _, _ = Create_File(str(user))

        # 保存PPT文件（保留/规范扩展名，不使用temp后缀）
        orig_ext = os.path.splitext(ppt_file.filename)[1].lower()
        if orig_ext not in ['.ppt', '.pptx', '.pptm']:
            orig_ext = '.pptx'
        ppt_filename = f"{user}{orig_ext}"
        ppt_final_path = os.path.join(save_user_path, ppt_filename)
        ppt_file.save(ppt_final_path)

        # 解析PPT备注
        try:
            # 尝试导入python-pptx库
            import pptx
            from pptx import Presentation

            # 打开PPT文件（添加多种错误处理）
            prs = None
            ppt_remakes = {}

            try:
                prs = Presentation(ppt_final_path)
            except Exception as e:
                logger.error(f"PPT文件格式问题: {e}")
                # 尝试使用win32com作为备用方案
                try:
                    from util.PPT2Video import Ppt_2_Video
                    ppt_processor = Ppt_2_Video("", "", "")
                    if ppt_processor.Open_Ppt(ppt_final_path):
                        # 获取幻灯片数量
                        slide_count = ppt_processor.presentation.Slides.Count
                        logger.info(f"使用COM接口打开PPT，共{slide_count}页")

                        # 生成默认备注
                        for i in range(1, slide_count + 1):
                            ppt_remakes[f"Slide {i}"] = f"这是第{i}页的内容"

                        ppt_processor.Close()

                        # 保存解析结果
                        ppt_remake_filename = Save_PPT_Remake(save_user_path, ppt_remakes)

                        if ppt_remakes:
                            logger.info(f"使用COM接口成功解析PPT文件，共{len(ppt_remakes)}页")
                            return jsonify(result="Success", data={
                                "parsed_count": len(ppt_remakes),
                                "remakes": ppt_remakes,
                                "method": "com_fallback"
                            })

                except Exception as com_e:
                    logger.error(f"COM接口也失败: {com_e}")

                return jsonify(result="Failed", error="PPT文件格式不支持或已损坏，请尝试重新保存PPT文件")
            except Exception as e:
                logger.error(f"PPT解析失败: {e}")
                # 如果python-pptx失败，也尝试使用COM接口
                try:
                    from util.PPT2Video import Ppt_2_Video
                    ppt_processor = Ppt_2_Video("", "", "")
                    if ppt_processor.Open_Ppt(ppt_final_path):
                        slide_count = ppt_processor.presentation.Slides.Count
                        logger.info(f"使用COM接口打开PPT，共{slide_count}页")

                        for i in range(1, slide_count + 1):
                            ppt_remakes[f"Slide {i}"] = f"这是第{i}页的内容"

                        ppt_processor.Close()

                        ppt_remake_filename = Save_PPT_Remake(save_user_path, ppt_remakes)

                        if ppt_remakes:
                            logger.info(f"使用COM接口成功解析PPT文件，共{len(ppt_remakes)}页")
                            return jsonify(result="Success", data={
                                "parsed_count": len(ppt_remakes),
                                "remakes": ppt_remakes,
                                "method": "com_fallback"
                            })

                except Exception as com_e:
                    logger.error(f"COM接口也失败: {com_e}")

                return jsonify(result="Failed", error=f"PPT解析失败: {str(e)}")

            # 如果成功打开，使用python-pptx解析备注（增强版：notes_text_frame 和 shapes 文本双通道）
            if prs:
                def _extract_notes_from_slide(_slide):
                    chunks = []
                    try:
                        if _slide.has_notes_slide:
                            ns = _slide.notes_slide
                            # 优先使用标准 notes_text_frame
                            tf = getattr(ns, 'notes_text_frame', None)
                            if tf is not None and getattr(tf, 'paragraphs', None):
                                for p in tf.paragraphs:
                                    try:
                                        t = (getattr(p, 'text', None) or '').strip()
                                        if t:
                                            chunks.append(t)
                                    except Exception:
                                        continue
                            # 回退：遍历 notes 页中的 shapes 文本
                            if not chunks and hasattr(ns, 'shapes'):
                                for sh in ns.shapes:
                                    try:
                                        if getattr(sh, 'has_text_frame', False) and getattr(sh, 'text', None):
                                            t = (sh.text or '').strip()
                                            if t:
                                                chunks.append(t)
                                    except Exception:
                                        continue
                    except Exception as e2:
                        logger.warning(f"解析备注文本失败（回退后）: {e2}")
                    text = "\n".join(chunks).strip()
                    return text

                # 遍历每一页幻灯片
                for slide_idx, slide in enumerate(prs.slides, 1):
                    try:
                        notes_text = _extract_notes_from_slide(slide)
                    except Exception as notes_e:
                        logger.warning(f"获取第{slide_idx}页备注失败: {notes_e}")
                        notes_text = ""

                    if notes_text:
                        ppt_remakes[f"Slide {slide_idx}"] = notes_text
                    else:
                        ppt_remakes[f"Slide {slide_idx}"] = f"这是第{slide_idx}页的内容"

                # 保存解析结果
                ppt_remake_filename = Save_PPT_Remake(save_user_path, ppt_remakes)

                # 保留PPT文件供后续生成PPT视频使用
                logger.info(f"PPT文件已保存为: {ppt_final_path}")

                if ppt_remakes:
                    logger.info(f"成功解析PPT文件，共{len(ppt_remakes)}页备注")
                    return jsonify(result="Success", data={
                        "parsed_count": len(ppt_remakes),
                        "remakes": ppt_remakes,
                        "method": "python_pptx"
                    })
                else:
                    return jsonify(result="Failed", error="PPT文件中没有找到备注内容")

        except ImportError:
            # 如果没有安装python-pptx库，直接使用COM接口
            logger.warning("python-pptx库未安装，尝试使用COM接口")
            try:
                from util.PPT2Video import Ppt_2_Video
                ppt_processor = Ppt_2_Video("", "", "")
                if ppt_processor.Open_Ppt(ppt_final_path):
                    slide_count = ppt_processor.presentation.Slides.Count
                    logger.info(f"使用COM接口打开PPT，共{slide_count}页")

                    ppt_remakes = {}
                    for i in range(1, slide_count + 1):
                        ppt_remakes[f"Slide {i}"] = f"这是第{i}页的内容"

                    ppt_processor.Close()

                    ppt_remake_filename = Save_PPT_Remake(save_user_path, ppt_remakes)

                    if ppt_remakes:
                        logger.info(f"使用COM接口成功解析PPT文件，共{len(ppt_remakes)}页")
                        return jsonify(result="Success", data={
                            "parsed_count": len(ppt_remakes),
                            "remakes": ppt_remakes,
                            "method": "com_only"
                        })

                return jsonify(result="Failed", error="PPT解析失败，无法打开PPT文件")

            except Exception as e:
                logger.error(f"COM接口解析失败: {e}")
                return jsonify(result="Failed", error=f"系统未安装PPT解析库，且COM接口失败: {str(e)}")

        except Exception as e:
            logger.error(f"PPT解析异常: {e}")
            return jsonify(result="Failed", error=f"PPT解析异常: {str(e)}")

    except Exception as e:
        logger.error(f"上传PPT文件失败: {e}")
        return jsonify(result="Failed", error=f"处理PPT文件失败: {str(e)}")

# 根据音频时长生成PPT视频
@app.route('/Generate_PPT_Video', methods=['POST'])
def Generate_PPT_Video():
    """根据音频时长数据生成对应的PPT视频"""
    try:
        POST_JSON = request.get_json()
        user = POST_JSON.get('User')

        if not user:
            return jsonify(result="Failed", error="缺少用户信息")

        _, _, _, save_user_path, _, _ = Create_File(str(user))

        # 检查Time.json文件是否存在
        time_json_path = os.path.join(save_user_path, "Time.json")
        if not os.path.exists(time_json_path):
            return jsonify(result="Failed", error="找不到时间数据文件")

        # 读取时间数据
        with open(time_json_path, 'r', encoding='utf-8') as f:
            time_data = json.load(f)

        if not time_data:
            return jsonify(result="Failed", error="时间数据为空")

        logger.info(f"开始为用户 {user} 生成PPT视频，共 {len(time_data)} 个片段")

        # 导入PPT处理模块
        try:
            from util.PPT2Video import Ppt_2_Video
            import pptx
            from pptx import Presentation
        except ImportError as e:
            logger.error(f"导入PPT处理模块失败: {e}")
            return jsonify(result="Failed", error="系统缺少PPT处理模块")

        # 检查是否存在原始PPT文件（先找非临时文件，再找临时文件）
        ppt_files = [f for f in os.listdir(save_user_path) if f.endswith(('.ppt', '.pptx')) and 'temp' not in f]
        if not ppt_files:
            # 如果没有找到非临时文件，查找临时文件
            ppt_files = [f for f in os.listdir(save_user_path) if f.endswith(('.ppt', '.pptx')) and 'temp' in f]
            if not ppt_files:
                return jsonify(result="Failed", error="找不到PPT文件，请先上传PPT文件")
            logger.info(f"使用临时PPT文件: {ppt_files}")
        else:
            logger.info(f"使用正式PPT文件: {ppt_files}")

        # 使用最新的PPT文件
        ppt_file = sorted(ppt_files)[-1]
        ppt_path = os.path.join(save_user_path, ppt_file)
        logger.info(f"Resolve PPT path => rel: {ppt_path}, abs: {os.path.abspath(ppt_path)}, exists: {os.path.exists(ppt_path)}, cwd: {os.getcwd()}")

        # 为本次PPT视频生成使用独立的临时工作目录，避免占用/删除用户实际的 Mov_Video 目录
        temp_base = os.path.join(save_user_path, "__ppt_temp")
        output_frames = os.path.join(temp_base, "output_frames")
        input_frames = os.path.join(temp_base, "input_frames")
        mov_temp = os.path.join(temp_base, "mov_video")

        for folder in [output_frames, input_frames, mov_temp]:
            if os.path.exists(folder):
                try:
                    shutil.rmtree(folder)
                except Exception as e:
                    logger.warning(f"清理临时目录失败(忽略继续): {folder} -> {e}")
            os.makedirs(folder, exist_ok=True)

        # 初始化PPT处理对象
        ppt_processor = Ppt_2_Video(output_frames, input_frames, mov_temp)

        # 打开PPT文件（优先COM，后续可用CreateVideo直接导出）
        try:
            opened_by_com = ppt_processor.Open_Ppt(ppt_path)
            if not opened_by_com:
                logger.error("COM 打开PPT失败，将尝试使用 python-pptx 进行回退导出")
        except Exception as e:
            opened_by_com = False
            logger.error(f"打开PPT文件失败: {e} | path={ppt_path} | abs={os.path.abspath(ppt_path)} | exists={os.path.exists(ppt_path)}")

        # 根据音频时长生成视频片段
        slide_indices = sorted(time_data.keys(), key=lambda x: int(x.split()[-1]) if x.split() else float('inf'))
        logger.info(f"时长数据键已排序: {slide_indices}")
        total_time = 0

        for slide_idx in slide_indices:
            if slide_idx not in time_data:
                continue

            audio_duration = time_data[slide_idx]
            total_time += audio_duration

            logger.info(f"处理幻灯片 {slide_idx}，音频时长: {audio_duration}秒")

            # 设置PPT幻灯片切换时间
            try:
                slide_index = int(slide_idx.split()[-1]) - 1  # "Slide 1" -> 0
                ppt_processor.Set_Ppt_Transtion_Speed(slide_index, audio_duration)
            except (ValueError, IndexError) as e:
                logger.warning(f"设置幻灯片时间失败: {e}")

        # 生成基础PPT视频：若COM已成功打开则优先使用CreateVideo；否则回退到python-pptx+moviepy
        try:
            # 直接写入目标目录（先写临时文件，后原子替换）
            user_name = os.path.basename(save_user_path)
            user_result_ppt_path = os.path.normpath(os.path.join('Result', 'ppt', user_name))
            os.makedirs(user_result_ppt_path, exist_ok=True)
            final_ppt_video_path = os.path.join(user_result_ppt_path, "PPT_Video.mp4")
            tmp_ppt_video_path = os.path.join(user_result_ppt_path, "__PPT_Video.tmp.mp4")
            logger.info(f"开始生成PPT基础视频(临时): {tmp_ppt_video_path}")

            # 若 COM 已打开，优先使用 PowerPoint 的 CreateVideo 导出
            if opened_by_com:
                logger.info("使用 COM CreateVideo 导出PPT视频")
                # 通过 util.PPT2Video 的 Create_Ppt_Base_Video 导出
                ok = ppt_processor.Create_Ppt_Base_Video(tmp_ppt_video_path)
                if not ok:
                    logger.warning("COM CreateVideo 导出失败，回退到 python-pptx + moviepy")

            # 回退路径：使用 python-pptx + moviepy 逐页导出图片再合成
            if (not os.path.exists(tmp_ppt_video_path)) or (os.path.getsize(tmp_ppt_video_path) <= 0):
                try:
                    from moviepy.editor import ImageClip, concatenate_videoclips
                    # 创建临时目录保存每一页的图片
                    temp_dir = os.path.join(save_user_path, "__ppt_temp", "temp_slides")
                    os.makedirs(temp_dir, exist_ok=True)

                    # 使用python-pptx重新打开PPT文件来导出图片
                    try:
                        prs = Presentation(ppt_path)
                    except Exception as e:
                        logger.error(f"打开PPT文件失败: {e}")
                        # 如果python-pptx失败，跳过图片生成，直接返回成功
                        return jsonify(result="Success", data={
                            "video_path": "PPT_Video.mp4",
                            "total_duration": total_time,
                            "method": "no_images"
                        })

                    # 保存每一页为图片
                    for slide_idx, slide in enumerate(prs.slides):
                        slide_img_path = os.path.join(temp_dir, f"slide_{slide_idx}.png")

                        # 导出幻灯片为图片
                        slide_img = None
                        if slide.shapes:
                            # 寻找第一个图片对象，跳过占位符
                            for shape in slide.shapes:
                                if hasattr(shape, 'image'):
                                    slide_img = shape.image
                                    break

                        if slide_img:
                            # 如果有图片，直接保存
                            with open(slide_img_path, 'wb') as f:
                                f.write(slide_img.blob)
                        else:
                            # 如果没有图片，创建空白图片
                            from PIL import Image, ImageDraw, ImageFont
                            img = Image.new('RGB', (1920, 1080), color='white')
                            draw = ImageDraw.Draw(img)

                            # 添加幻灯片编号
                            draw.text((50, 50), f"Slide {slide_idx}", fill='black')

                            # 如果有备注，添加备注内容
                            if slide.has_notes_slide:
                                notes_slide = slide.notes_slide
                                notes_text = ""
                                for paragraph in notes_slide.notes_text_frame.paragraphs:
                                    notes_text += paragraph.text.strip() + "\n"
                                if notes_text.strip():
                                    # 分行显示备注
                                    lines = notes_text.strip().split('\n')
                                    for i, line in enumerate(lines[:5]):  # 最多显示5行
                                        draw.text((50, 150 + i * 50), line.strip(), fill='black')

                            img.save(slide_img_path)

                    # 使用时间数据设置每张图片的显示时长
                    image_clips = []
                    slide_duration = {}  # 记录每张幻灯片的时长

                    # 从time_data获取时长信息
                    for slide_idx, slide in enumerate(prs.slides):
                        slide_key = f"Slide {slide_idx + 1}"
                        if slide_key in time_data:
                            slide_duration[slide_idx] = time_data[slide_key]
                        else:
                            slide_duration[slide_idx] = 5.0  # 默认5秒

                    # 创建图片剪辑
                    for slide_idx in range(len(prs.slides)):
                        slide_img_path = os.path.join(temp_dir, f"slide_{slide_idx}.png")
                        if os.path.exists(slide_img_path):
                            duration = slide_duration.get(slide_idx, 5.0)
                            clip = ImageClip(slide_img_path).set_duration(duration)
                            image_clips.append(clip)

                    if image_clips:
                        # 合成视频
                        final_clip = concatenate_videoclips(image_clips)
                        final_clip.write_videofile(tmp_ppt_video_path, fps=24, codec='libx264', audio=False)
                        logger.info(f"PPT基础视频生成成功(临时): {tmp_ppt_video_path}")
                    else:
                        # 如果没有图片，创建空白视频
                        create_blank_video(tmp_ppt_video_path, 10.0)  # 10秒空白视频
                        logger.warning("没有找到有效的幻灯片内容，创建空白视频")

                except ImportError:
                    logger.error("缺少moviepy库，无法生成PPT视频")
                    return jsonify(result="Failed", error="系统缺少视频处理库")
                except Exception as e:
                    logger.error(f"生成PPT视频过程出错: {e}")
                    return jsonify(result="Failed", error=f"PPT视频生成过程出错: {str(e)}")
                finally:
                    # 清理临时目录
                    try:
                        if os.path.exists(temp_dir):
                            shutil.rmtree(temp_dir)
                            logger.info(f"清理临时目录: {temp_dir}")
                    except Exception as e:
                        logger.warning(f"清理临时目录失败(忽略继续): {temp_dir} -> {e}")

            # 临时文件原子替换为正式文件
            if os.path.exists(tmp_ppt_video_path):
                try:
                    os.replace(tmp_ppt_video_path, final_ppt_video_path)
                except Exception:
                    # 某些文件系统上 os.replace 可能失败，退回为 copy + remove
                    shutil.copy2(tmp_ppt_video_path, final_ppt_video_path)
                    try:
                        os.remove(tmp_ppt_video_path)
                    except Exception:
                        pass
                logger.info(f"PPT视频生成成功(发布): {final_ppt_video_path}")
                try:
                    size_bytes = os.path.getsize(final_ppt_video_path)
                except Exception:
                    size_bytes = -1
                logger.info(f"PPT视频文件是否存在: {os.path.exists(final_ppt_video_path)}, size: {size_bytes}")

                return jsonify(result="Success", data={
                    "video_path": os.path.relpath(final_ppt_video_path, start=base_dir).replace("\\", "/"),
                    "total_duration": total_time
                })
            else:
                return jsonify(result="Failed", error="PPT视频临时文件生成失败")

        except Exception as e:
            logger.error(f"生成PPT视频失败: {e}")
            return jsonify(result="Failed", error=f"PPT视频生成失败: {str(e)}")
        finally:
            # 尝试清理本次PPT生成的临时工作区
            try:
                if os.path.exists(temp_base):
                    shutil.rmtree(temp_base)
                    logger.info(f"清理PPT生成临时工作区: {temp_base}")
            except Exception as e:
                logger.warning(f"清理PPT生成临时工作区失败(忽略): {temp_base} -> {e}")
            # 确保关闭 COM 的 PowerPoint
            try:
                if opened_by_com:
                    ppt_processor.Close()
            except Exception as e:
                logger.warning(f"关闭PowerPoint失败: {e}")

    except Exception as e:
        logger.error(f"生成PPT视频接口错误: {e}")
        return jsonify(result="Failed", error=f"生成PPT视频失败: {str(e)}")

# 保存PPT备注信息
@app.route('/Send_PPT_Remakes', methods=['POST'])
def Set_PPT_Remakes():
    POST_JSON = request.get_json()
    user = POST_JSON.get('User')
    ppt_remakes = POST_JSON.get("PPT_Remakes")

    try:
        _, _, _, save_user_path, _, _ = Create_File(str(user))
        # 兼容前端把 PPT_Remakes 当作字符串传入的情况
        if isinstance(ppt_remakes, str):
            try:
                parsed = json.loads(ppt_remakes)
                ppt_remakes = parsed
            except Exception as pe:
                logger.warning(f"PPT_Remakes 解析为 JSON 失败，按原样保存字符串: {pe}")
        ppt_remake_filename = Save_PPT_Remake(save_user_path,ppt_remakes)

        with open(ppt_remake_filename, 'r', encoding='utf-8') as f:
            data = json.load(f)

        if(data != {}):
            return jsonify(result="Success")
        else:
            return jsonify(result="Failed")
    except:
        return jsonify(result="Failed")
    
#保存真人照片
@app.route('/Send_Image', methods=['POST'])
def Send_Image():
    # POST_JSON = request.get_json()
    # user = POST_JSON.get('User')
    # img = request.files.get('Image')  # 从post请求中获取图片数据
    
    POST_JSON = request.get_json()
    user = POST_JSON.get('User')
    img_data_base64 = POST_JSON.get('Img')
    
    try:
        _, _, _, save_user_path, _, _ = Create_File(str(user))
        
        # 解码base64字符串并保存到文件
        img_data = base64.b64decode(img_data_base64)
        with open("img.png", "wb") as img_file:
            img_file.write(img_data)
        Save_Image(save_user_path,"img.png")
        
        return jsonify(result="Success")
    except:
        return jsonify(result="Failed")

#保存教师视频
@app.route('/Send_Teacher_Video', methods=['POST'])
def Send_Teacher_Video():
    string = request.form.get('Json')
    video = request.files['File'].read()
    try:
        json_data = json.loads(string)
        user = json_data.get('User')
        _, _, _, save_user_path, _, _ = Create_File(str(user))
        
        # 写入文件
        with open(os.path.join(save_user_path, "Video.mp4"), "wb") as video_file:
            video_file.write(video)
        
        # 处理视频文件的保存逻辑
        # Save_Video(save_user_path, "video.mp4")
        
        return jsonify(result="Success")
    except Exception as e:
        print(e)
        return jsonify(result="Failed")
    
#获取VITS音频时长
@app.route('/Recive_Wav_Time', methods=['POST'])
def Recive_Wav_Time():
    POST_JSON = request.get_json()
    user = POST_JSON.get('User')

    try:
        result_vits_user_path, _, _, save_user_path, _, _ = Create_File(str(user))
        wav_time_dict = Save_Tiem(save_user_path, result_vits_user_path)
        return  jsonify(result = wav_time_dict)
        
    except:
        return jsonify(result="Failed")
 
# 获取用户音频时长
@app.route('/Recive_User_Wav_Time', methods=['POST'])
def Recive_User_Wav_Time():
    POST_JSON = request.get_json()
    user = POST_JSON.get('User')

    try:
        _, _, _, save_user_path, _, _ = Create_File(str(user))
        ppt_audio_dir = os.path.join(save_user_path, "PPT_Audio")
        audio_json_save_path = os.path.join(save_user_path, "Audio_save_path.json")
        user_wav_path = os.path.join(save_user_path, "PPT_Audio")
        
        Write_Json(ppt_audio_dir, audio_json_save_path)
        wav_time_dict = Save_Tiem(save_user_path, user_wav_path)
        return  jsonify(result = wav_time_dict)
        
    except:
        return jsonify(result="Failed")
 
# 只读调试端点：检查用户当前产物与关键路径
@app.route('/Debug_Pipeline', methods=['POST'])
def Debug_Pipeline():
    try:
        POST_JSON = request.get_json()
        user = POST_JSON.get('User') if POST_JSON else None
        if not user:
            return jsonify(result="Failed", error="缺少用户信息")

        result_vits_user_path, result_sadtalker_user_path, result_wav2lip_user_path, save_user_path, _, _ = Create_File(str(user))

        # 基本路径
        user_name = os.path.basename(save_user_path)
        user_result_ppt_path = os.path.normpath(os.path.join('Result', 'ppt', user_name))
        final_ppt_video_path = os.path.join(user_result_ppt_path, "PPT_Video.mp4")

        # 查找PPT文件
        try:
            ppt_files = [f for f in os.listdir(save_user_path) if f.endswith(('.ppt', '.pptx'))]
        except Exception:
            ppt_files = []
        ppt_files_sorted = sorted(ppt_files)
        latest_ppt = ppt_files_sorted[-1] if ppt_files_sorted else None
        latest_ppt_abs = os.path.abspath(os.path.join(save_user_path, latest_ppt)) if latest_ppt else None

        # 读取关键配置/产物
        def safe_json_load(path):
            try:
                with open(path, 'r', encoding='utf-8') as f:
                    return json.load(f)
            except Exception:
                return None

        audio_json = safe_json_load(os.path.join(save_user_path, "Audio_save_path.json"))
        video_json = safe_json_load(os.path.join(save_user_path, "Video_save_path.json"))
        time_json = safe_json_load(os.path.join(save_user_path, "Time.json"))
        ppt_remake_json = safe_json_load(os.path.join(save_user_path, "PPT_Remake.json"))
        people_loc_json = safe_json_load(os.path.join(save_user_path, "People_Location.json"))

        def file_status(path):
            if not path:
                return {"exists": False}
            try:
                return {
                    "exists": os.path.exists(path),
                    "abs": os.path.abspath(path),
                    "size": os.path.getsize(path) if os.path.exists(path) else -1
                }
            except Exception:
                return {"exists": False, "abs": os.path.abspath(path)}

        resp = {
            "paths": {
                "save_user_path": os.path.abspath(save_user_path),
                "result_vits_path": os.path.abspath(result_vits_user_path),
                "result_sadtalker_path": os.path.abspath(result_sadtalker_user_path),
                "result_wav2lip_path": os.path.abspath(result_wav2lip_user_path),
            },
            "ppt": {
                "candidates": ppt_files_sorted,
                "latest": latest_ppt,
                "latest_abs": latest_ppt_abs,
                "latest_exists": os.path.exists(latest_ppt_abs) if latest_ppt_abs else False
            },
            "manifests": {
                "Audio_save_path": audio_json,
                "Video_save_path": video_json,
                "Time": time_json,
                "PPT_Remake": ppt_remake_json,
                "People_Location": people_loc_json
            },
            "products": {
                "ppt_video": file_status(final_ppt_video_path)
            }
        }

        # 汇总音频/视频文件数量
        try:
            vits_files = [f for f in os.listdir(result_vits_user_path) if f.lower().endswith('.wav')]
        except Exception:
            vits_files = []
        try:
            sadtalker_files = [f for f in os.listdir(result_sadtalker_user_path) if f.lower().endswith(('.mp4', '.mov'))]
        except Exception:
            sadtalker_files = []

        resp["counts"] = {
            "vits_wav_count": len(vits_files),
            "sadtalker_video_count": len(sadtalker_files)
        }

        # 关键检查：Slide 1 是否在 manifests 中
        def first_missing(manifest):
            if isinstance(manifest, dict):
                for k, v in sorted(manifest.items(), key=lambda x: (str(x[0]))):
                    if not v or not os.path.exists(v):
                        return {"key": k, "path": v, "exists": os.path.exists(v) if v else False}
            return None

        resp["sanity"] = {
            "audio_first_missing": first_missing(audio_json) if isinstance(audio_json, dict) else "Audio manifest not dict",
            "video_first_missing": first_missing(video_json) if isinstance(video_json, dict) else "Video manifest not dict",
        }

        return jsonify(result="Success", data=resp)
    except Exception as e:
        logger.error(f"Debug_Pipeline 错误: {e}")
        return jsonify(result="Failed", error=str(e))

#接收前端视频
@app.route('/Send_Video', methods=['POST'])
def Send_Video():
    string = request.form.get('Json')
    video = request.files['File'].read()

    try:
        json_data = json.loads(string)
        user = json_data.get('User')
        _, _, _, save_user_path, _, _ = Create_File(str(user))

        # 使用正确的PPT视频保存路径
        user_name = os.path.basename(save_user_path)
        user_result_ppt_path = os.path.normpath(os.path.join('Result', 'ppt', user_name))
        final_ppt_video_path = os.path.join(user_result_ppt_path, "PPT_Video.mp4")

        # 确保PPT视频目录存在
        os.makedirs(user_result_ppt_path, exist_ok=True)

        # 写入文件
        with open(final_ppt_video_path, "wb") as video_file:
            video_file.write(video)

        logger.info(f"用户上传的视频已保存到: {final_ppt_video_path}")
        logger.info(f"视频文件是否存在: {os.path.exists(final_ppt_video_path)}")

        return jsonify(result="Success")
    except Exception as e:
        logger.error(f"保存视频失败: {e}")
        return jsonify(result="Failed")

# 保存数字人插入页数的json
@app.route('/Send_People_Location', methods=['POST'])
def Send_People_Location():
    POST_JSON = request.get_json()
    user = POST_JSON.get('User')
    people_location = POST_JSON.get("People_Location")

    try:
        _, _, _, save_user_path, _, _ = Create_File(str(user))
        people_location_filename = Save_People_Location(save_user_path,people_location)
        
        with open(people_location_filename, 'r', encoding='utf-8') as f:
            data = json.load(f)
            
        if(data != {}):
            return jsonify(result="Success")
        else:
            return jsonify(result="Failed")
    except:
        return jsonify(result="Failed")   

# 保存用于插入PPT的音频
@app.route('/Send_PPT_Audio', methods=['POST'])
def Send_PPT_Audio():
    string = request.form.get('Json')
    audio = request.files['File'].read()
    
    try:
        json_data = json.loads(string)
        user = json_data.get('User')
        audio_name = json_data.get('Audio_Name')
        
        _, _, _, save_user_path, _, _ = Create_File(str(user))
        Save_Insert_Audio(save_user_path, audio_name, audio)
        
        return jsonify(result="Success")
    except:
        return jsonify(result="Failed")

#####################################################################################
#                                 配置参数                                           #
#####################################################################################

# 配置所有模型参数
@app.route('/Send_Config', methods=['POST'])
def Send_Config():
    POST_JSON = request.get_json()
    user = POST_JSON.get('User')
    vits_config = POST_JSON.get('VITS_Config')
    sadtalker_config = POST_JSON.get('SadTalker_Config')
    
    try:
        _, _, _, save_user_path, _, _ = Create_File(str(user))
        
        Config_SadTalker_Parmes(save_user_path, sadtalker_config)
        Config_VITS_Parmes(save_user_path, vits_config)
        
        return jsonify(result="Success")
    except:
        return jsonify(result="Failed")
    
    
#配置wav2lip参数
@app.route('/Send_Wav2Lip_Config', methods=['POST'])
def Send_Wav2Lip_Config():
    POST_JSON = request.get_json()
    user = POST_JSON.get('User')
    wav2lip_config = POST_JSON.get('Wav2Lip_Config')
    
    try:
        _, _, _, save_user_path, _, _ = Create_File(str(user))
        
        Config_Wav2Lip_Parmes(save_user_path, wav2lip_config)
        
        return jsonify(result="Success")
    except:
        return jsonify(result="Failed")
    
#选择训练的VITS模型
@app.route('/Send_Select_Train_VITS_Model', methods=['POST'])
def Send_Select_Train_VITS_Model():
    POST_JSON = request.get_json()
    user = POST_JSON.get('User')
    
    try:
        _, _, _, save_user_path, _, _ = Create_File(str(user))
        Select_Train_VITS_Model(save_user_path,user)
        
        return jsonify(result="Success")
    
    except:
        return jsonify(result="Failed")

#选择VITS模型
@app.route('/Send_Select_VITS_Model', methods=['POST'])
def Send_Select_VITS_Model():
    POST_JSON = request.get_json()
    user = POST_JSON.get('User')
    index = POST_JSON.get('Index')
    
    try:
        _, _, _, save_user_path, _, _ = Create_File(str(user))
        
        Select_VITS_Model(save_user_path,str(index))
        
        return jsonify(result="Success")
    except:
        return jsonify(result="Failed")
   
    
#####################################################################################
#                                VITS功能                                           #
#####################################################################################

#保存用于训练VITS的音频
@app.route('/Send_Tarin_Audio', methods=['POST'])
def Send_Tarin_Audio():
    string = request.form.get('Json')
    audio = request.files['File'].read()
    
    try:
        json_data = json.loads(string)
        user = json_data.get('User')
        audio_name = json_data.get('Audio_Name')
        
        _, _, _, save_user_path, _, _ = Create_File(str(user))
        Save_Train_Audio(save_user_path, audio_name, audio)
        
        return jsonify(result="Success")
    except:
        return jsonify(result="Failed")

#训练VITS模型
@app.route('/Train_VITS_Model', methods=['POST'])
def Train_VITS_Model():
    POST_JSON = request.get_json()
    user = POST_JSON.get('User')
    json_data = POST_JSON.get('Label')

    try:
        _, _, _, save_user_path, _, _ = Create_File(str(user))
        Task_State(save_user_path, "VITS_Train", "Processing")
        executor.submit(Train_VITS, save_user_path, user, json_data)

        return jsonify(result="VITS_Train")
    except:
        return jsonify(result="Failed")
    
#保存VITS的参照音频跟文字
@app.route('/Send_Ref_Wav_And_Text', methods=['POST'])
def Send_Ref_Wav_And_Text():
    string = request.form.get('Json')
    audio = request.files['File'].read()
    
    try:
        json_data = json.loads(string)
        user = json_data.get('User')
        ref_text = (json_data.get('Ref_Text') or '').strip()
        _, _, _, save_user_path, _, _ = Create_File(user)
        
        file =  save_user_path + "/" + "Audio.mp3"
         # 写入文件
        with open(file, "wb") as audio_file:
            audio_file.write(audio)

        # 若未提供或提供的是默认占位文本，则尝试进行轻量ASR以生成参考文本
        def _is_placeholder(t: str) -> bool:
            if not t:
                return True
            s = t.strip()
            return s in ("这是参考文本", "", "请输入参考文本")

        transcript = ref_text
        if _is_placeholder(ref_text):
            try:
                # 优先使用 faster-whisper（更快），其次回退 openai/whisper
                model = None
                try:
                    from faster_whisper import WhisperModel  # type: ignore
                    model = WhisperModel("small", device="cpu", compute_type="int8")
                    segments, info = model.transcribe(file, language="zh")
                    transcript = "".join([seg.text for seg in segments]).strip()
                except Exception:
                    import warnings; warnings.filterwarnings("ignore")
                    try:
                        import whisper  # type: ignore
                        wmodel = whisper.load_model("small")
                        res = wmodel.transcribe(file, language="zh")
                        transcript = (res.get("text") or "").strip()
                    except Exception:
                        pass
                if not transcript:
                    transcript = ""
            except Exception as asr_e:
                # ASR失败不阻断流程
                transcript = ref_text

        Save_VITS_Ref_Wav_And_Text(save_user_path, file,  {"Text" : transcript}, "None")
            
        return jsonify(result="Success")
    except Exception as e:
        print(e)
        return jsonify(result="Failed")
    
#获取训练后的VITS模型的名字
@app.route('/Get_Train_VITS_Model_Name', methods=['POST'])
def Get_Train_VITS_Model_Name():
    POST_JSON = request.get_json()
    user = POST_JSON.get('User')
    
    try:
        _, _, _, save_user_path, _, _ = Create_File(str(user))
        weightPath = os.path.join(save_user_path, "Weight")
        #判断文件夹不为空
        if len(os.listdir(weightPath)) > 0:
            return jsonify(result=f"{user}")
        else:
            return jsonify(result="Failed")
        
    except:
        return jsonify(result="Failed")
    
#####################################################################################
#                                      推理                                         #
#####################################################################################
   
#推理效果展示视频
@app.route('/Get_Test_Inference', methods=['POST'])
@log_exceptions()
def Get_Test_Inference():
    POST_JSON = request.get_json() or {}
    user = POST_JSON.get('User')
    result_vits_user_path, result_sadtalker_user_path, _, save_user_path, _, _ = Create_File(user)
    DH = VITS_Sadtalker_Join(result_vits_user_path, result_sadtalker_user_path, save_user_path)
    # Get_Parmes 返回 (sad, vits, wav2lip)，这里之前少接收一个导致 ValueError
    sad_parames_yaml_path, vits_parames_yaml_path, _ = Get_Parmes(save_user_path)
    DH.Set_Params_and_Model(sad_parames_yaml_path, vits_parames_yaml_path)

    with open(os.path.join(save_user_path,"Ref_text.json"), "r", encoding='utf-8') as f:
        data = json.load(f)

    ref_wav_path = os.path.join(save_user_path,"Ref_Wav.wav")
    ref_text = data["Text"]
    test = "你好，我是数字人授课录制系统，很高兴为您服务。"
    video_output_path = os.path.join(save_user_path,"Test")
    imag_path = os.path.join(save_user_path,"Image.png")

    # 1) 生成测试音频
    audio_path = DH.Inference_VITS_test(ref_wav_path, ref_text, test)
    # 2) 生成测试视频
    DH.Inference_SadTalker_test(imag_path, audio_path, video_output_path)

    # 3) 将本次测试的音/视频写入 JSON（供后续可选合成 PPT 使用）
    try:
        test_video_path = os.path.join(save_user_path, "Test.mp4")
        audio_json = os.path.join(save_user_path, "Audio_save_path.json")
        video_json = os.path.join(save_user_path, "Video_save_path.json")
        time_json = os.path.join(save_user_path, "Time.json")

        with open(audio_json, 'w', encoding='utf-8') as jf:
            json.dump({"0": audio_path}, jf, ensure_ascii=False)
        with open(video_json, 'w', encoding='utf-8') as jf:
            json.dump({"0": test_video_path}, jf, ensure_ascii=False)
        # 计算音频时长并写入 Time.json
        try:
            with wave.open(audio_path, 'rb') as wavf:
                frames = wavf.getnframes(); fr = wavf.getframerate()
                duration = frames / float(fr) if fr else 0.0
        except Exception as te:
            logger.warning(f"计算测试音频时长失败: {te}")
            duration = 0.0
        with open(time_json, 'w', encoding='utf-8') as jf:
            json.dump({"0": duration}, jf, ensure_ascii=False)
    except Exception as e:
        logger.warning(f"写入测试 JSON 失败(不影响直接返回测试视频): {e}")

    # 4) 若请求显式要求或环境变量开启，且已存在 PPT_Video.mp4，则自动进行 PPT 合成
    auto_merge_flag = POST_JSON.get('MergePPT')
    if isinstance(auto_merge_flag, str):
        auto_merge = auto_merge_flag.lower() in ("1","true","yes")
    else:
        auto_merge = bool(auto_merge_flag)
    if not auto_merge:
        auto_merge = os.environ.get("APP_TEST_AUTOMERGE_PPT", "0") in ("1","true","True")

    ppt_video = os.path.join(save_user_path, "PPT_Video.mp4")
    if auto_merge and os.path.exists(ppt_video):
        try:
            # 同步执行完整合成流程，返回最终合成视频
            from Main import Remove_Video_Background, Video_Joint, Last_Video_Join_Audio
            Remove_Video_Background(save_user_path)
            merged_mid = Video_Joint(save_user_path)
            final_video = Last_Video_Join_Audio(save_user_path, merged_mid)
            video_data_base64 = Encode_Video(final_video)
            return jsonify(result=video_data_base64, merged=True)
        except Exception as me:
            logger.exception(f"Test 模式自动贴合 PPT 失败，将返回测试视频: {me}")

    # 默认返回测试视频
    video_data_base64 = Encode_Video(os.path.join(save_user_path, "Test.mp4"))
    return jsonify(result=video_data_base64, merged=False)
       
#推理VITS(单个)
@app.route('/Get_Inference_VITS', methods=['POST'])
@log_exceptions()
def Get_Inference_VITS():
    POST_JSON = request.get_json() or {}
    user = POST_JSON.get('User')
    result_vits_user_path, result_sadtalker_user_path, _, save_user_path, _, _ = Create_File(user)
    DH = VITS_Sadtalker_Join(result_vits_user_path, result_sadtalker_user_path, save_user_path)
    sad_parames_yaml_path, vits_parames_yaml_path, _ = Get_Parmes(save_user_path)
    DH.Set_Params_and_Model(sad_parames_yaml_path, vits_parames_yaml_path)

    with open(os.path.join(save_user_path,"Ref_text.json"), "r", encoding='utf-8') as f:
        data = json.load(f)

    ref_wav_path = os.path.join(save_user_path,"Ref_Wav.wav")
    ref_text = data["Text"]
    test = "你好，我是数字人授课录制系统，很高兴为您服务。"

    DH.Inference_VITS_test(ref_wav_path, ref_text, test)
    return jsonify(result="Success")
      
#推理VITS(多个)
@app.route('/Get_Inference_VITS_Multiple', methods=['POST'])
@log_exceptions()
def Get_Inference_VITS_Multiple():
    POST_JSON = request.get_json() or {}
    user = POST_JSON.get('User')
    result_vits_user_path, result_sadtalker_user_path, _, save_user_path, _, _ = Create_File(user)
    Task_State(save_user_path, "VITS_Inference", "Processing")
    VITS_Multiple_Inference(result_vits_user_path, result_sadtalker_user_path, save_user_path)
    return jsonify(result="VITS_Inference")
      
# 推理VITS跟Sadtalker
@app.route('/Get_Inference_VITS_Sadtalker', methods=['POST'])
@log_exceptions()
def Get_Inference_VITS_Sadtalker():
    POST_JSON = request.get_json() or {}
    user = POST_JSON.get('User')
    result_vits_user_path, result_sadtalker_user_path, _, save_user_path, _, _ = Create_File(user)
    Task_State(save_user_path, "Audio_Video_Inference", "Processing")
    executor.submit(VITS_Sadtalker_Inference, result_vits_user_path, result_sadtalker_user_path, save_user_path)
    return jsonify(result="Audio_Video_Inference")

# 推理用户音频跟Sadtalker
@app.route('/Get_Inference_User_Audio_Sadtalker', methods=['POST'])
@log_exceptions()
def Get_Inference_User_Audio_Sadtalker():
    POST_JSON = request.get_json() or {}
    user = POST_JSON.get('User')
    _, result_sadtalker_user_path, _, save_user_path, _, _ = Create_File(user)
    Task_State(save_user_path, "Audio_Video_Inference", "Processing")
    executor.submit(User_Wav_Sadtalker_Inference, result_sadtalker_user_path, save_user_path)
    return jsonify(result="Audio_Video_Inference")
    
    
# 推理VITS跟Wav2Lip
@app.route('/Get_Inference_VITS_Wav2Lip', methods=['POST'])
@log_exceptions()
def Get_Inference_VITS_Wav2Lip():
    POST_JSON = request.get_json() or {}
    user = POST_JSON.get('User')
    result_vits_user_path, _, result_wav2lip_user_path, save_user_path, _, _ = Create_File(user)
    Task_State(save_user_path, "Audio_Video_Inference", "Processing")
    VITS_Wav2Lip_Inference(result_vits_user_path, result_wav2lip_user_path, save_user_path)
    return jsonify(result="Audio_Video_Inference")

# 推理用户音频跟Wav2Lip
@app.route('/Get_Inference_User_Audio_Wav2Lip', methods=['POST'])
@log_exceptions()
def Get_Inference_User_Audio_Wav2Lip():
    POST_JSON = request.get_json() or {}
    user = POST_JSON.get('User')
    _, _, result_wav2lip_user_path, save_user_path, _, _ = Create_File(user)
    Task_State(save_user_path, "Audio_Video_Inference", "Processing")
    executor.submit(User_Wav_Wav2Lip_Inference, result_wav2lip_user_path, save_user_path)
    return jsonify(result="Audio_Video_Inference")

# ppt跟视频合成（合成最终效果视频，全插入数字人）
@app.route('/PPT_Video_Merge', methods=['POST'])
def PPT_Video_Merge():
    POST_JSON = request.get_json()
    user = POST_JSON.get('User')
    try:
        _, _, _, save_user_path, _, _ = Create_File(user)
        # 将任务提交到线程池
        Task_State(save_user_path, "Video_Merge", "Processing")
        executor.submit(Video_Merge,save_user_path)

        return jsonify(result = "Video_Merge")
    except:
        return jsonify(result="Failed")

# ppt跟视频合成（合成最终效果视频，可选择插入数字人）
@app.route('/PPT_Video_Merge_Select_Into', methods=['POST'])
def PPT_Video_Merge_Select_Into():
    POST_JSON = request.get_json()
    user = POST_JSON.get('User')
    try:
        _, _, _, save_user_path, _, _ = Create_File(user)
                # 将任务提交到线程池
        Task_State(save_user_path, "Video_Merge", "Processing")
        executor.submit(Video_Merge_Select_Into,save_user_path)

        return jsonify(result = "Video_Merge")
    except:
        return jsonify(result="Failed")
        
# ppt跟视频合成（没有数字人）
@app.route('/PPT_Video_Merge_No_Into', methods=['POST'])
def PPT_Video_Merge_No_Into():
    POST_JSON = request.get_json()
    user = POST_JSON.get('User')
    try:
        _, _, _, save_user_path, _, _ = Create_File(user)
        Video_Join_Audio(save_user_path)
        
        return jsonify(result = "Success")
    except:
        return jsonify(result="Failed")

#####################################################################################
#                                    获取                                           #
#####################################################################################

# 拉取视频
@app.route('/Pull_Video_Merge', methods=['POST'])
def Pull_Video_merge():
    POST_JSON = request.get_json()
    user = POST_JSON.get('User')
    try:
        _, _, _, _, save_user_path, user_result_ppt_path, user_result_generation_path = Create_File(user)
        last_video = os.path.join(user_result_generation_path, "last_video.mp4")
        video_data_base64 = Encode_Video(last_video)
        return jsonify(result = video_data_base64)
    except:
        return jsonify(result="Failed")

# 直接流式返回最终合成视频（浏览器/播放器友好）
@app.route('/Download_Merged_Video', methods=['GET'])
def Download_Merged_Video():
    try:
        user = request.args.get('User') or request.args.get('user')
        filename = request.args.get('filename')
        logger.info(f"[Download_Merged_Video] 请求参数: user={user}, filename={filename}")
        if not user:
            logger.error("[Download_Merged_Video] 缺少用户信息")
            return jsonify(result="Failed", error="缺少用户信息"), 400

        user_result_generation_path = os.path.normpath(os.path.join('Result','generation',str(user)))
        logger.info(f"[Download_Merged_Video] 用户视频目录: {user_result_generation_path}")

        if filename:
            final_path = os.path.join(user_result_generation_path, filename)
        else:
            final_path = os.path.join(user_result_generation_path, 'last_video.mp4')

        logger.info(f"[Download_Merged_Video] 视频文件路径: {final_path}")
        if not os.path.exists(final_path):
            logger.error(f"[Download_Merged_Video] 未找到指定视频文件: {final_path}")
            return jsonify(result="Failed", error="未找到指定视频文件"), 404

        file_size = os.path.getsize(final_path)
        logger.info(f"[Download_Merged_Video] 视频文件存在, size={file_size}")
        if file_size == 0:
            logger.error(f"[Download_Merged_Video] 视频文件大小为0字节: {final_path}")

        download_name = filename if filename else 'last_video.mp4'
        try:
            resp = send_file(final_path, mimetype='video/mp4', as_attachment=False, download_name=download_name)
        except Exception as send_e:
            logger.error(f"[Download_Merged_Video] send_file异常: {send_e}")
            return jsonify(result="Failed", error=f"send_file异常: {str(send_e)}"), 500

        # 禁止缓存，暴露必要头部
        resp.headers['Cache-Control'] = 'no-store'
        resp.headers['Access-Control-Expose-Headers'] = 'Content-Length, Content-Range'
        logger.info(f"[Download_Merged_Video] 响应已准备，文件名: {download_name}, size: {file_size}")
        return resp
    except Exception as e:
        logger.exception(f"[Download_Merged_Video] 失败: {e}")
        return jsonify(result="Failed", error=str(e)), 500
 
 
# 拉取推理的VITS声音
@app.route('/Pull_VITS_Audio', methods=['POST'])
def Pull_VITS_Audio():
    POST_JSON = request.get_json()
    user = POST_JSON.get('User')
    try:
        _, _, _, save_user_path, _, _ = Create_File(user)
        vits_wav = os.path.join(save_user_path,"Test_VITS.wav")
        wav_data_base64 = Encode_Video(vits_wav)
        return jsonify(result=wav_data_base64)
    except:
        return jsonify(result="Failed")

# 获取用户生成的视频列表
@app.route('/Get_Video_List', methods=['GET'])
def Get_Video_List():
    """获取用户生成的视频列表"""
    try:
        user = request.args.get('User')
        if not user:
            return jsonify(result="Failed", error="缺少用户参数")

        # 约定最终视频输出位置：Data/generation/<user>/
        user_result_generation_path = os.path.normpath(os.path.join('Result','generation',str(user)))
        os.makedirs(user_result_generation_path, exist_ok=True)

        videos = []

        # 查找所有视频文件
        if os.path.exists(user_result_generation_path):
            for filename in os.listdir(user_result_generation_path):
                if filename.lower().endswith(('.mp4', '.mov', '.avi', '.mkv')):
                    filepath = os.path.join(user_result_generation_path, filename)
                    try:
                        # 获取文件信息
                        stat = os.stat(filepath)
                        file_size = stat.st_size
                        create_time = stat.st_ctime

                        # 生成视频URL（用于下载和播放）
                        video_url = f"/Download_Merged_Video?User={user}&filename={filename}"

                        # 获取视频时长（简单估算，实际可以用moviepy获取精确时长）
                        # 这里使用文件大小粗略估算时长（假设5MB ≈ 1分钟）
                        estimated_duration = max(1, int(file_size / (5 * 1024 * 1024) * 60))

                        videos.append({
                            "id": filename.replace('.', '_'),  # 用于前端标识
                            "name": filename,
                            "url": video_url,
                            "duration": estimated_duration,
                            "size": file_size,
                            "createTime": create_time
                        })
                    except Exception as e:
                        logger.warning(f"处理视频文件失败 {filepath}: {e}")
                        continue

        # 按创建时间倒序排列
        videos.sort(key=lambda x: x['createTime'], reverse=True)

        return jsonify(result=videos)

    except Exception as e:
        logger.error(f"获取视频列表失败: {e}")
        return jsonify(result="Failed", error=str(e))

# 删除用户视频
@app.route('/Delete_Video', methods=['POST'])
def Delete_Video():
    """删除用户生成的视频文件"""
    try:
        POST_JSON = request.get_json()
        user = POST_JSON.get('User')
        video_id = POST_JSON.get('VideoId')

        if not user or not video_id:
            return jsonify(result="Failed", error="缺少必要参数")

        # 将id转回文件名
        filename = video_id.replace('_', '.')
        user_result_generation_path = os.path.normpath(os.path.join('Result', 'generation', str(user)))
        video_path = os.path.join(user_result_generation_path, filename)

        if not os.path.exists(video_path):
            return jsonify(result="Failed", error="视频文件不存在")

        try:
            os.remove(video_path)
            logger.info(f"成功删除视频文件: {video_path}")
            return jsonify(result="Success")
        except Exception as e:
            logger.error(f"删除视频文件失败: {e}")
            return jsonify(result="Failed", error=f"删除失败: {str(e)}")

    except Exception as e:
        logger.error(f"删除视频接口错误: {e}")
        return jsonify(result="Failed", error=str(e))
 
 
    
if __name__ == '__main__':
    
    app.run("0.0.0.0")


